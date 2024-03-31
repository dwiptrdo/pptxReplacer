from fastapi.responses import FileResponse
from python_pptx_text_replacer import TextReplacer
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE 
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor      
from pptx.util import Cm,Pt  
from pptx.util import Inches
from pptx import Presentation
from fastapi import FastAPI
from pptx.parts.image import Image as PptImage
from PIL import Image
import base64
import datetime
import pptx
import io



app = FastAPI()


# Fungsi untuk mengganti data dalam chart pada slide tertentu
def replace_chart_with_data(slide, chart_index, chart_data):
    chart_count = 0
    for shape in slide.shapes:
        if shape.has_chart:
            chart_count += 1
            if chart_count == chart_index + 1:  
                chart = shape.chart
                chart.replace_data(chart_data)
                print(f"Chart with index {chart_index} found and replaced successfully on the specified slide.")
                return
    print(f"Chart with index {chart_index} not found on the specified slide.")


# fungsi untuk mengkonversi warna yang ada di dalam data
def hex_to_rgb(hex_color):
    # Hapus tanda pagar (#) jika ada
    hex_color = hex_color.lstrip('#')
    # Konversi string heksadesimal menjadi nilai RGB
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return (r, g, b)



# Endpoint untuk generate
@app.post("/generate")
def generate(file: dict):
    
    prs = Presentation('template/file/20240117_setneg_biweekly.pptx')

    ## menambahkan shape
    data = file['result']
    percentage1 = data['sna']['clusters'][0]['percentage']
    label1 = data['sna']['clusters'][0]['label']

    percentage2 = data['sna']['clusters'][1]['percentage']
    label2 = data['sna']['clusters'][1]['label']

    percentage3 = data['sna']['clusters'][2]['percentage']
    label3 = data['sna']['clusters'][2]['label']

    nomor_slide = 3
    slide = prs.slides[nomor_slide]

    # untuk mengatur warna yang sesuai untuk shape
    color1 = data['sna']['clusters'][0]['color']
    color2 = data['sna']['clusters'][1]['color']
    color3 = data['sna']['clusters'][2]['color']

    rgb_color1 = hex_to_rgb(color1)
    rgb_color2 = hex_to_rgb(color2)
    rgb_color3 = hex_to_rgb(color3)

    # menambahkan shape ke slide
    # shape 1
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,    # Jenis shape (rounded rectangle)
        Cm(0.72), Cm(9.56),             # Koordinat x, y (dalam sentimeter)
        Cm(5), Cm(1.3))                 # Lebar dan tinggi (dalam sentimeter)

    # Atur warna isi bentuk
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*rgb_color2) # mengatur warna shape

    # Atur teks di dalam shape
    text_frame = rect.text_frame
    p = text_frame.paragraphs[0]
    p.text = f'{label2} ({percentage2:.1f}%)'  # Teks yang ingin ditampilkan
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(11)


    # shape 2
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,    
        Cm(25.33), Cm(1.8),             
        Cm(5), Cm(1.3))                 

    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*rgb_color1) 

    text_frame = rect.text_frame
    p = text_frame.paragraphs[0]
    p.text = f'{label1} ({percentage1:.1f}%)' 
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(11)

    # shape 3
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,   
        Cm(25.33), Cm(14.3),           
        Cm(5), Cm(1.3))                

    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*rgb_color3) 

    text_frame = rect.text_frame
    p = text_frame.paragraphs[0]
    p.text = f'{label3} ({percentage3:.1f}%)' 
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(11)


    ## menambahkan text box
    summary3 = data['sna']['clusters'][2]
    summary_text = summary3['summary'][0]
    slide = prs.slides[3] 


    # Tentukan posisi dan ukuran kotak teks
    left_inch = Inches(10)
    top_inch = Inches(5.88)
    width_inch = Inches(3.25)
    height_inch = Inches(2)


    # Tambahkan text box ke slide
    text_box = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    text_frame = text_box.text_frame


    # Tambahkan teks ke dalam text box
    p = text_frame.add_paragraph()
    p.text = summary_text
    # mengatur size dan word wrap
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.word_wrap = True

    # Atur warna teks
    for run in p.runs:
        run.font.color.rgb = RGBColor(255, 255, 255)

    # Atur ukuran font teks
    p.font.size = Pt(11)


    ## menghapus shape yang tidak di perlukan
    slide_index = 3 
    slide = prs.slides[slide_index]
    shapes_to_delete = [9, 8, 7, 6]  

    # Hapus bentuk (shape) dari slide
    for shape_index in shapes_to_delete:
        slide.shapes[shape_index]._element.getparent().remove(slide.shapes[shape_index]._element)


    ## mereplace chart
    data = file['result']['each_day_count']
    dates = []
    counts = []

    for item in data:
        for date, count in item.items():
            dates.append(date)
            counts.append(count)

    
    dates = [datetime.datetime.strptime(date, "%Y-%m-%d").date() for date in dates]

    
    chart_data = CategoryChartData()
    chart_data.categories = dates
    chart_data.add_series('', counts)


    slide_index = 1
    chart_index_to_replace = 0

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break


    
    image_percakapan1 = file['result']['list_of_images'][0]

    binary_dataper1 = base64.b64decode(image_percakapan1)

    image = Image.open(io.BytesIO(binary_dataper1))

    ## menyimpan gambar dalam format PNG
    save_percakapan1 = ('assets/percakapan1.png')
    image.save(save_percakapan1)
    image.close()

    


    image_percakapan2 = file['result']['list_of_images'][1]

    binary_dataper2 = base64.b64decode(image_percakapan2)

    image = Image.open(io.BytesIO(binary_dataper2))

    save_percakapan2 = ('assets/percakapan2.png')
    image.save(save_percakapan2)
    image.close()

    

    image_percakapan3 = file['result']['list_of_images'][2]

    binary_dataper3 = base64.b64decode(image_percakapan3)

    image = Image.open(io.BytesIO(binary_dataper3))

    save_percakapan3 = ('assets/percakapan3.png')
    image.save(save_percakapan3)
    image.close()


    

    no_image = ('assets/noimage.jpg')

    new_image_paths = [save_percakapan3, save_percakapan2, save_percakapan1, no_image, no_image, no_image]  

    
    img_shape_indexes = [9, 7, 6, 4, 5, 8]  # Ganti dengan indeks gambar yang diinginkan di slide tersebut

    # menentukan indeks slide yang ingin diperbarui
    slide_index = 2 

    # Iterasi melalui daftar gambar baru dan indeks gambar di slide
    for new_image_path, img_shape_index in zip(new_image_paths, img_shape_indexes):
        
        new_pptx_img = PptImage.from_file(new_image_path)

        img_shape = prs.slides[slide_index].shapes[img_shape_index]

        slide_part, rId = img_shape.part, img_shape._element.blip_rId
        image_part = slide_part.related_part(rId)

        image_part.blob = new_pptx_img.blob




    image_polarisasi = file['result']['sna']['image']

    binary_datapolarisasi = base64.b64decode(image_polarisasi)

    # Membuka gambar menggunakan PIL (Pillow)
    image = Image.open(io.BytesIO(binary_datapolarisasi))

    # Simpan gambar dalam format PNG
    save_polarisasi = ('assets/polarisasi.png')
    image.save(save_polarisasi)
    image.close()


    polarisasi_image = save_polarisasi

    new_pptx_img = pptx.parts.image.Image.from_file(polarisasi_image)

    # menentukan index slide pptx dan index gambar
    img_shape = prs.slides[3].shapes[0]  

    slide_part, rId = img_shape.part, img_shape._element.blip_rId
    image_part = slide_part.related_part(rId)

    image_part.blob = new_pptx_img._blob


    save_file = ('template/20240117_setneg_biweekly.pptx')
    prs.save(save_file)


    data = file['result']
    data_ikn = file['result']['topic']

    total_twitter = data["platform_count"][0]["twitter"]["total"]
    percentage_twitter = data["platform_count"][0]["twitter"]["percentage"]


    total_fb = data["platform_count"][3]["facebook"]["total"]
    percentage_fb = data["platform_count"][3]["facebook"]["percentage"]


    total_youtube = data["platform_count"][1]["youtube"]["total"]
    percentage_youtube = data["platform_count"][1]["youtube"]["percentage"]


    total_instagram = data["platform_count"][2]["instagram"]["total"]
    percentage_instagram = data["platform_count"][2]["instagram"]["percentage"]


    total_tiktok = data["platform_count"][4]["tiktok"]["total"]
    percentage_tiktok = data["platform_count"][4]["tiktok"]["percentage"]


    summary1 = data['sna']['clusters'][0]
    summary2 = data['sna']['clusters'][1]

    replacer = TextReplacer(save_file, slides='', tables=True, charts=False, textframes=True)

    replacer.replace_text([
        ('IKN', data_ikn),
        ('1 – 15 Januari 2024', data['earliest_date'] + ' Sampai ' + data['latest_date']),
        ('1 - 15 Januari 2024', data['earliest_date'] + ' Sampai ' + data['latest_date']),
        ('36.343', str(data['total_count'])),
        ('Perhatian terhadap IKN cenderung tidak meningkat namun sempat mengalami lonjakan karena isu-isu tertentu.', data['trend_analysis']),
        ('Pasca debat, Prabowo dikritik karena kepemilikan tanah di IKN.', data['topics'][0]),
        ('Pemerintah dikritik karena ingin menggelontorkan triliunan uang untuk pembangunan IKN.', data['topics'][1]),
        ('Netizen soroti isu Djarum dan Wings Group hengkang dari konsorsium IKN.', data['topics'][2]),
        ('7.961', str(data['sna']['statistics']['account_count'])),
        ('100', str(data['sna']['statistics']['hashtag_count'])),
        ('22.542', str(data['sna']['statistics']['activity_count'])),
        ('Kelompok pro IKN aktif angkat keberhasilan Jokowi mendapat investasi 7 triliun setelah mengadakan kunjungan ke ASEAN, terutama dari Brunei.', data['sna']['summary'][0]),
        ('Kelompok kontra IKN masih terus menyindir IKN karena menunjukkan sikap pemerintah yang hanya ingin menguntungkan diri sendiri, bukan pro rakyat.', data['sna']['summary'][1]),
        ('Kontra IKN juga angkat isu adanya investor IKN yang keluar dari konsorsium.', data['sna']['summary'][2]),
        ('26.951 Data (74%)', f'{total_twitter} Data ({percentage_twitter * 100}%)'),
        ('586 Data (1.6%)', f'{total_fb} Data ({percentage_fb * 100}%)'),
        ('7.351 Data (20%)', f'{total_youtube} Data ({percentage_youtube * 100}%)'),
        ('726 Data (1.99%)', f'{total_instagram} Data ({percentage_instagram * 100}%)'),
        ('729 Data (2%)', f'{total_tiktok} Data ({percentage_tiktok * 100}%)'),
        ('Kelompok ini cenderung merupakan akun-akun pro pemerintah dan pro Prabowo.', summary2['summary'][0]),
        ('Kelompok ini angkat keberhasilan pemerintah mendapat investasi 7 Triliun setelah kunjungan ke negara-negara Asean.', summary2['summary'][1]),
        ('Kelompok pro Prabowo kritik Anies yang dianggap sikapnya kini mulai tidak konsisten terhadap IKN yang sebelumnya aktif menolak.', summary2['summary'][2]),
        ('Kelompok ini klarifikasi isu adanya investor yang mundur dari IKN.', summary2['summary'][3]),
        ('Kelompok ini menunjukkan dampak positif IKN terhadap daerah sekitarnya yang akan ikut maju.', ''),
        ('Kelompok kontra cenderung lebih banyak dari pendukung Anies.', summary1['summary'][0]),
        ('Kelompok ini terus mengkritik proyek IKN yang disebut terus menyedot uang negara hingga triliunan, sehingga anggaran tidak digunakan untuk memperhatikan rakyat.', summary1['summary'][1]),
        ('Jokowi dikritik “menjadi jubir Prabowo” untuk membela masalah kritik netizen terhadap kesejahteraan TNI dibandingkan dengan pembangunan IKN.', summary1['summary'][2]),
        ('Kelompok ini viralkan adanya penghentian pengerjaan pipa air oleh pemilik lahan di IKN karena mengaku belum mendapat bayaran dari Pemerintah.', summary1['summary'][3]),
        ('Pendukung Anies secara khusus menyindir Gibran dan Jokowi yang sebelumnya menyebut sudah memiliki investor, namun ternyata ditinggalkan oleh Djarum dan Wings Group.', summary1['summary'][4]),
        ('Muncul sikap bila Jakarta punya masalah, selesaikan, bukan pindah ke IKN.', summary1['summary'][5])
    ])

    file_output = "result/20240117_setneg_biweekly.pptx"
    replacer.write_presentation_to_file(file_output)

    return FileResponse(file_output, filename="20240117_setneg_biweekly.pptx")